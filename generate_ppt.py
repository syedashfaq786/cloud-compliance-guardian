"""
Invecto Compliance Guard — Innovation Presentation Generator
Generates a professional dark-themed PowerPoint matching the application's UI/UX.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Colors matching the app UI ─────────────────────────────────────────────
BG_DARK = RGBColor(0x0D, 0x12, 0x24)       # #0d1224
BG_CARD = RGBColor(0x1A, 0x1D, 0x2E)       # #1a1d2e
ACCENT = RGBColor(0xFF, 0x9F, 0x43)         # #ff9f43 (orange)
ACCENT_DIM = RGBColor(0xFF, 0x7A, 0x00)     # #ff7a00
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_SEC = RGBColor(0x94, 0xA3, 0xB8)       # #94a3b8
GREEN = RGBColor(0x22, 0xC5, 0x5E)
RED = RGBColor(0xEF, 0x44, 0x44)
PURPLE = RGBColor(0x8B, 0x5C, 0xF6)
BLUE = RGBColor(0x4F, 0x6E, 0xF7)
CYAN = RGBColor(0x38, 0xBD, 0xF8)

# ── Paths ──────────────────────────────────────────────────────────────────
BASE = os.path.dirname(os.path.abspath(__file__))
LOGO_DIR = os.path.join(BASE, "dashboard", "public", "logos")
SCREENSHOT_DIR = r"C:\Users\ashfa\.gemini\antigravity\brain\d5887426-1a81-4b86-adcf-3a4a32ae31ff"

SCREENSHOTS = {
    "login": os.path.join(SCREENSHOT_DIR, "screenshot_login.png"),
    "dashboard": os.path.join(SCREENSHOT_DIR, "screenshot_dashboard.png"),
    "monitoring": os.path.join(SCREENSHOT_DIR, "screenshot_monitoring.png"),
    "topology": os.path.join(SCREENSHOT_DIR, "screenshot_topology.png"),
    "compliance": os.path.join(SCREENSHOT_DIR, "screenshot_compliance_rules.png"),
}

OUTPUT = os.path.join(BASE, "Invecto_Compliance_Guard_Presentation.pptx")


def set_slide_bg(slide, color=BG_DARK):
    """Set slide background to a solid color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rounded_rect(slide, left, top, width, height, fill_color=BG_CARD, border_color=None):
    """Add a rounded rectangle shape."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=14, color=WHITE,
                 bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    """Add a text box with styled text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_text(slide, left, top, width, height, items, font_size=13, color=WHITE,
                    bullet_color=ACCENT, spacing=Pt(6)):
    """Add bulleted text list."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = spacing
        
        # Add orange bullet
        run_bullet = p.add_run()
        run_bullet.text = "▸ "
        run_bullet.font.size = Pt(font_size)
        run_bullet.font.color.rgb = bullet_color
        run_bullet.font.name = "Calibri"

        # Check for bold prefix (text before colon)
        if ":" in item and not item.startswith("✅") and not item.startswith("🔄") and not item.startswith("🎯"):
            parts = item.split(":", 1)
            run_bold = p.add_run()
            run_bold.text = parts[0] + ":"
            run_bold.font.size = Pt(font_size)
            run_bold.font.color.rgb = ACCENT
            run_bold.font.bold = True
            run_bold.font.name = "Calibri"

            run_rest = p.add_run()
            run_rest.text = parts[1]
            run_rest.font.size = Pt(font_size)
            run_rest.font.color.rgb = color
            run_rest.font.name = "Calibri"
        else:
            run_text = p.add_run()
            run_text.text = item
            run_text.font.size = Pt(font_size)
            run_text.font.color.rgb = color
            run_text.font.name = "Calibri"

    return txBox


def add_section_header(slide, left, top, width, text, color=ACCENT, size=20):
    """Add a section header with underline accent."""
    add_text_box(slide, left, top, width, Inches(0.5), text, font_size=size, color=color, bold=True)
    # Orange underline
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top + Inches(0.42), Inches(1.2), Pt(3))
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()


def add_slide_number(slide, num, total=8):
    """Add slide number to bottom right."""
    add_text_box(slide, Inches(11.5), Inches(7.1), Inches(1.5), Inches(0.3),
                 f"{num} / {total}", font_size=10, color=TEXT_SEC, alignment=PP_ALIGN.RIGHT)


def add_logo_badge(slide, left, top):
    """Add the Invecto brand badge."""
    # Icon box
    rect = add_rounded_rect(slide, left, top, Inches(0.45), Inches(0.45),
                            fill_color=RGBColor(0x2A, 0x1F, 0x0D), border_color=ACCENT)
    # Try to add actual logo
    logo_path = os.path.join(LOGO_DIR, "aws.png")  # We'll use text instead
    
    # Brand text
    add_text_box(slide, left + Inches(0.55), top - Inches(0.02), Inches(2), Inches(0.28),
                 "Invecto", font_size=14, color=WHITE, bold=True)
    add_text_box(slide, left + Inches(0.55), top + Inches(0.22), Inches(2), Inches(0.22),
                 "COMPLIANCE GUARD", font_size=9, color=ACCENT, bold=True)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE BUILDERS
# ═══════════════════════════════════════════════════════════════════════════════

def build_slide_1(prs):
    """Title & Team Overview"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    set_slide_bg(slide)

    # Decorative gradient bar at top
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.33), Pt(4))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()

    add_logo_badge(slide, Inches(0.6), Inches(0.6))

    # Main title
    add_text_box(slide, Inches(0.6), Inches(1.6), Inches(10), Inches(0.7),
                 "Invecto Compliance Guard", font_size=40, color=WHITE, bold=True)

    add_text_box(slide, Inches(0.6), Inches(2.3), Inches(10), Inches(0.5),
                 "AI-Powered Cloud Compliance Platform", font_size=22, color=ACCENT, bold=True)

    # Problem statement card
    card = add_rounded_rect(slide, Inches(0.6), Inches(3.2), Inches(11.5), Inches(1.3),
                            fill_color=BG_CARD, border_color=RGBColor(0x2A, 0x2D, 0x3E))

    add_text_box(slide, Inches(0.9), Inches(3.35), Inches(2), Inches(0.3),
                 "THE PROBLEM", font_size=11, color=ACCENT, bold=True)

    add_text_box(slide, Inches(0.9), Inches(3.7), Inches(11), Inches(0.7),
                 '"Cloud misconfigurations cause 65% of security incidents — enterprises lack a unified, '
                 'AI-driven compliance engine that audits across multiple frameworks, cloud providers, '
                 'and container platforms in real time."',
                 font_size=15, color=TEXT_SEC)

    # Team card
    card2 = add_rounded_rect(slide, Inches(0.6), Inches(5.0), Inches(5.5), Inches(1.8),
                             fill_color=BG_CARD, border_color=RGBColor(0x2A, 0x2D, 0x3E))

    add_text_box(slide, Inches(0.9), Inches(5.15), Inches(2), Inches(0.3),
                 "TEAM", font_size=11, color=ACCENT, bold=True)

    add_bullet_text(slide, Inches(0.9), Inches(5.5), Inches(5), Inches(1.2), [
        "Syed Ashfaq Mohiddin — Project Lead",
        "Swastik Shetty — Mentor",
    ], font_size=15)

    # Tech badges
    badges = ["AWS", "Azure", "GCP", "Terraform", "Docker", "K8s", "CIS", "NIST", "CCM"]
    for i, badge in enumerate(badges):
        x = Inches(6.5) + Inches(i % 3) * Inches(1.5)
        y = Inches(5.3) + Inches(i // 3) * Inches(0.5)
        rect = add_rounded_rect(slide, x, y, Inches(1.3), Inches(0.35),
                                fill_color=RGBColor(0x1A, 0x1A, 0x0D),
                                border_color=RGBColor(0x3D, 0x2E, 0x10))
        add_text_box(slide, x, y + Inches(0.04), Inches(1.3), Inches(0.3),
                     badge, font_size=11, color=ACCENT, bold=True, alignment=PP_ALIGN.CENTER)

    add_slide_number(slide, 1)


def build_slide_2(prs):
    """Problem Statement"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_logo_badge(slide, Inches(0.6), Inches(0.4))

    add_text_box(slide, Inches(0.6), Inches(1.1), Inches(10), Inches(0.5),
                 'The "Cyber-AI" Problem Statement', font_size=28, color=WHITE, bold=True)

    # The Gap
    card1 = add_rounded_rect(slide, Inches(0.6), Inches(1.9), Inches(11.5), Inches(1.5),
                             fill_color=BG_CARD, border_color=RGBColor(0x2A, 0x2D, 0x3E))
    add_text_box(slide, Inches(0.9), Inches(2.0), Inches(2), Inches(0.3),
                 "THE GAP", font_size=11, color=RED, bold=True)
    add_text_box(slide, Inches(0.9), Inches(2.35), Inches(11), Inches(1.0),
                 "Traditional compliance tools operate in silos — separate tools for AWS vs Azure vs GCP, "
                 "separate scanners for Terraform vs Docker vs Kubernetes, and manual mapping to CIS/NIST/CCM "
                 "frameworks. Security teams spend 70%+ of their time on manual audit configuration rather "
                 "than actual remediation.",
                 font_size=14, color=TEXT_SEC)

    # Impact - 4 stat cards
    stats = [
        ("$4.45M", "Avg. breach cost\n(IBM 2023)", RED),
        ("82%", "Breaches involve\ncloud data", ACCENT),
        ("60%", "Orgs fail first\ncompliance audit", PURPLE),
        ("₹250 Cr", "DPDP Act max\npenalty", BLUE),
    ]
    for i, (val, label, color) in enumerate(stats):
        x = Inches(0.6) + i * Inches(3)
        card = add_rounded_rect(slide, x, Inches(3.7), Inches(2.7), Inches(1.3),
                                fill_color=BG_CARD, border_color=RGBColor(0x2A, 0x2D, 0x3E))
        # Accent left bar
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Inches(3.7), Pt(4), Inches(1.3))
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.fill.background()
        
        add_text_box(slide, x + Inches(0.2), Inches(3.85), Inches(2.5), Inches(0.5),
                     val, font_size=28, color=color, bold=True)
        add_text_box(slide, x + Inches(0.2), Inches(4.4), Inches(2.5), Inches(0.5),
                     label, font_size=12, color=TEXT_SEC)

    # Target Audience
    card2 = add_rounded_rect(slide, Inches(0.6), Inches(5.3), Inches(11.5), Inches(1.5),
                             fill_color=BG_CARD, border_color=RGBColor(0x2A, 0x2D, 0x3E))
    add_text_box(slide, Inches(0.9), Inches(5.4), Inches(3), Inches(0.3),
                 "TARGET AUDIENCE", font_size=11, color=GREEN, bold=True)
    add_bullet_text(slide, Inches(0.9), Inches(5.7), Inches(11), Inches(1.0), [
        "Enterprise Cloud Teams & DevSecOps — managing multi-cloud infrastructure",
        "Financial Services — RBI-regulated banks requiring continuous compliance",
        "Healthcare & Critical Infrastructure — HIPAA, DPDP Act requirements",
    ], font_size=13)

    add_slide_number(slide, 2)


def build_slide_3(prs):
    """Proposed AI Solution"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_logo_badge(slide, Inches(0.6), Inches(0.4))

    add_text_box(slide, Inches(0.6), Inches(1.1), Inches(10), Inches(0.5),
                 "Proposed AI Solution", font_size=28, color=WHITE, bold=True)

    # Solution Overview
    card1 = add_rounded_rect(slide, Inches(0.6), Inches(1.8), Inches(5.5), Inches(1.5),
                             fill_color=BG_CARD, border_color=RGBColor(0x2A, 0x2D, 0x3E))
    add_text_box(slide, Inches(0.9), Inches(1.9), Inches(3), Inches(0.3),
                 "SOLUTION OVERVIEW", font_size=11, color=ACCENT, bold=True)
    add_text_box(slide, Inches(0.9), Inches(2.25), Inches(5), Inches(1.0),
                 "Bespoke AI-Enhanced Platform — built from scratch with a custom compliance "
                 "engine + Cisco Sec-8B LLM inference for context-aware auditing. Not an "
                 "out-of-the-box wrapper.",
                 font_size=13, color=TEXT_SEC)

    # AI Secret Sauce
    card2 = add_rounded_rect(slide, Inches(6.4), Inches(1.8), Inches(5.7), Inches(1.5),
                             fill_color=BG_CARD, border_color=RGBColor(0x2A, 0x2D, 0x3E))
    add_text_box(slide, Inches(6.7), Inches(1.9), Inches(3), Inches(0.3),
                 'THE "AI SECRET SAUCE"', font_size=11, color=PURPLE, bold=True)
    add_bullet_text(slide, Inches(6.7), Inches(2.2), Inches(5.2), Inches(1.0), [
        "Cisco Sec-8B LLM — context-aware IaC analysis",
        "Predictive compliance scoring + trend analysis",
        "Auto-generated HCL/YAML fix snippets",
        "Single scan → CIS + NIST + CCM simultaneously",
    ], font_size=12, bullet_color=PURPLE)

    # Architecture Diagram
    add_text_box(slide, Inches(0.6), Inches(3.6), Inches(5), Inches(0.3),
                 "CORE ARCHITECTURE", font_size=11, color=CYAN, bold=True)

    # 3 connected boxes
    box_data = [
        ("DATA INPUT", "Cloud APIs (AWS/Azure/GCP)\nTerraform files\nDocker/K8s manifests\nGitHub repos", CYAN),
        ("AI PROCESSING", "Cisco Sec-8B LLM Inference\nCIS/NIST/CCM Rule Engine\nSeverity Classification\nTopology Mapping", ACCENT),
        ("SECURITY ACTION", "Compliance Dashboard\nPDF/CSV/JSON Reports\nAuto-Remediation\nDrift Alerts", GREEN),
    ]

    for i, (title, content, color) in enumerate(box_data):
        x = Inches(0.6) + i * Inches(4.2)
        card = add_rounded_rect(slide, x, Inches(4.0), Inches(3.8), Inches(2.8),
                                fill_color=BG_CARD, border_color=color)
        # Top accent bar
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Inches(4.0), Inches(3.8), Pt(3))
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.fill.background()

        add_text_box(slide, x + Inches(0.15), Inches(4.15), Inches(3.5), Inches(0.3),
                     title, font_size=13, color=color, bold=True, alignment=PP_ALIGN.CENTER)

        # Step number circle
        add_text_box(slide, x + Inches(1.5), Inches(4.5), Inches(0.8), Inches(0.4),
                     f"Step {i+1}", font_size=10, color=TEXT_SEC, alignment=PP_ALIGN.CENTER)

        add_text_box(slide, x + Inches(0.2), Inches(4.9), Inches(3.4), Inches(1.8),
                     content, font_size=12, color=TEXT_SEC)

        # Arrow connector (except last)
        if i < 2:
            arrow_x = x + Inches(3.85)
            add_text_box(slide, arrow_x, Inches(5.2), Inches(0.35), Inches(0.4),
                         "→", font_size=24, color=ACCENT, bold=True, alignment=PP_ALIGN.CENTER)

    add_slide_number(slide, 3)


def build_slide_4(prs):
    """Technical Innovation"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_logo_badge(slide, Inches(0.6), Inches(0.4))

    add_text_box(slide, Inches(0.6), Inches(1.1), Inches(10), Inches(0.5),
                 'Technical Innovation & "Persistence"', font_size=28, color=WHITE, bold=True)

    # Left - Innovation
    card1 = add_rounded_rect(slide, Inches(0.6), Inches(1.8), Inches(5.7), Inches(3.0),
                             fill_color=BG_CARD, border_color=RGBColor(0x2A, 0x2D, 0x3E))
    add_text_box(slide, Inches(0.9), Inches(1.9), Inches(5), Inches(0.3),
                 "HOW WE DIFFER FROM ZSCALER / PRISMA CLOUD", font_size=11, color=ACCENT, bold=True)
    add_bullet_text(slide, Inches(0.9), Inches(2.25), Inches(5.2), Inches(2.4), [
        "Unified multi-framework (CIS + NIST + CCM) in single scan",
        "AI-powered via Cisco Sec-8B — not pattern-matching",
        "Full container compliance (Docker + K8s) alongside cloud",
        "Real-time topology visualization with compliance heatmap",
        "Privacy-first: all inference on private infrastructure",
    ], font_size=12)

    # Right - Self-healing
    card2 = add_rounded_rect(slide, Inches(6.5), Inches(1.8), Inches(5.7), Inches(3.0),
                             fill_color=BG_CARD, border_color=RGBColor(0x2A, 0x2D, 0x3E))
    add_text_box(slide, Inches(6.8), Inches(1.9), Inches(5), Inches(0.3),
                 "SELF-HEALING / RESILIENCE", font_size=11, color=GREEN, bold=True)
    add_bullet_text(slide, Inches(6.8), Inches(2.25), Inches(5.2), Inches(2.4), [
        "Scan cache persistence — survives server restarts",
        "Encrypted credential vault (.data directory)",
        "Graceful degradation to deterministic rule engine",
        "WebSocket real-time drift detection alerts",
    ], font_size=12, bullet_color=GREEN)

    # Bottom - Tech Stack
    card3 = add_rounded_rect(slide, Inches(0.6), Inches(5.1), Inches(11.6), Inches(1.9),
                             fill_color=BG_CARD, border_color=RGBColor(0x2A, 0x2D, 0x3E))
    add_text_box(slide, Inches(0.9), Inches(5.2), Inches(3), Inches(0.3),
                 "TECH STACK", font_size=11, color=PURPLE, bold=True)

    stack_items = [
        ("Backend", "Python, FastAPI, Uvicorn, Cisco Sec-8B", ACCENT),
        ("Frontend", "Next.js 14, React, Glassmorphic UI", CYAN),
        ("Cloud", "Boto3, Azure SDK, Google Cloud SDK", GREEN),
        ("Reports", "ReportLab PDF, CSV, JSON", PURPLE),
    ]
    for i, (label, tech, color) in enumerate(stack_items):
        x = Inches(0.9) + i * Inches(2.9)
        add_text_box(slide, x, Inches(5.55), Inches(2.5), Inches(0.25),
                     label, font_size=11, color=color, bold=True)
        add_text_box(slide, x, Inches(5.8), Inches(2.7), Inches(0.8),
                     tech, font_size=11, color=TEXT_SEC)

    add_slide_number(slide, 4)


def build_slide_5(prs):
    """Business Impact"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_logo_badge(slide, Inches(0.6), Inches(0.4))

    add_text_box(slide, Inches(0.6), Inches(1.1), Inches(10), Inches(0.5),
                 "Business Impact & Compliance", font_size=28, color=WHITE, bold=True)

    # Efficiency Gains
    card1 = add_rounded_rect(slide, Inches(0.6), Inches(1.8), Inches(5.7), Inches(2.2),
                             fill_color=BG_CARD, border_color=RGBColor(0x2A, 0x2D, 0x3E))
    add_text_box(slide, Inches(0.9), Inches(1.9), Inches(3), Inches(0.3),
                 "EFFICIENCY GAINS", font_size=11, color=GREEN, bold=True)
    add_bullet_text(slide, Inches(0.9), Inches(2.2), Inches(5.2), Inches(1.7), [
        "MTTD reduced by 85% (4hrs → 35 seconds)",
        "500+ controls in single scan — replaces 3 separate tools",
        "MTTR reduced 60% — auto-generated fix snippets",
        "10+ hours saved per audit cycle — auto-reports",
    ], font_size=12, bullet_color=GREEN)

    # Regulatory
    card2 = add_rounded_rect(slide, Inches(6.5), Inches(1.8), Inches(5.7), Inches(2.2),
                             fill_color=BG_CARD, border_color=RGBColor(0x2A, 0x2D, 0x3E))
    add_text_box(slide, Inches(6.8), Inches(1.9), Inches(3), Inches(0.3),
                 "REGULATORY ALIGNMENT", font_size=11, color=BLUE, bold=True)
    add_bullet_text(slide, Inches(6.8), Inches(2.2), Inches(5.2), Inches(1.7), [
        "DPDP Act (India) — data classification & access auditing",
        "RBI Guidelines — IAM, encryption, key rotation checks",
        "CIS Benchmarks v4.1 — cloud hardening gold standard",
        "NIST 800-53 & CSA CCM v4.1 — federal & cloud governance",
    ], font_size=12, bullet_color=BLUE)

    # Scalability
    card3 = add_rounded_rect(slide, Inches(0.6), Inches(4.3), Inches(11.6), Inches(2.5),
                             fill_color=BG_CARD, border_color=RGBColor(0x2A, 0x2D, 0x3E))
    add_text_box(slide, Inches(0.9), Inches(4.4), Inches(3), Inches(0.3),
                 "SCALABILITY", font_size=11, color=ACCENT, bold=True)

    scale_items = [
        ("Multi-Cloud Native", "AWS, Azure, GCP\nfrom a single dashboard", "☁"),
        ("Auto-Discovery", "All regions & resources\nper cloud account", "🔍"),
        ("CI/CD Integration", "GitHub Actions scan\non every push", "⚡"),
        ("Container Scale", "Unlimited Docker/K8s\nfile scanning", "📦"),
    ]
    for i, (title, desc, icon) in enumerate(scale_items):
        x = Inches(0.9) + i * Inches(2.9)
        add_text_box(slide, x, Inches(4.8), Inches(2.5), Inches(0.3),
                     f"{icon}  {title}", font_size=14, color=ACCENT, bold=True)
        add_text_box(slide, x, Inches(5.2), Inches(2.5), Inches(0.8),
                     desc, font_size=12, color=TEXT_SEC)

    add_slide_number(slide, 5)


def build_slide_6(prs):
    """Roadmap"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_logo_badge(slide, Inches(0.6), Inches(0.4))

    add_text_box(slide, Inches(0.6), Inches(1.1), Inches(10), Inches(0.5),
                 "Roadmap & PoC", font_size=28, color=WHITE, bold=True)

    # Current Stage badge
    badge = add_rounded_rect(slide, Inches(0.6), Inches(1.8), Inches(4), Inches(0.5),
                             fill_color=RGBColor(0x0D, 0x2A, 0x15),
                             border_color=GREEN)
    add_text_box(slide, Inches(0.9), Inches(1.85), Inches(3.5), Inches(0.4),
                 "● CURRENT STAGE:  Working PoC — Fully Functional Prototype",
                 font_size=12, color=GREEN, bold=True)

    # Two-week plan
    card1 = add_rounded_rect(slide, Inches(0.6), Inches(2.6), Inches(5.7), Inches(1.6),
                             fill_color=BG_CARD, border_color=RGBColor(0x2A, 0x2D, 0x3E))
    add_text_box(slide, Inches(0.9), Inches(2.7), Inches(3), Inches(0.3),
                 'THE "TWO-WEEK PLAN"', font_size=11, color=ACCENT, bold=True)
    add_bullet_text(slide, Inches(0.9), Inches(3.0), Inches(5.2), Inches(1.1), [
        "Week 1: Azure + GCP live scan parity, drift detection",
        "Week 2: GitHub Actions CI/CD, scheduled scans, RBAC",
    ], font_size=13)

    # Milestones timeline
    card2 = add_rounded_rect(slide, Inches(0.6), Inches(4.5), Inches(11.6), Inches(2.5),
                             fill_color=BG_CARD, border_color=RGBColor(0x2A, 0x2D, 0x3E))
    add_text_box(slide, Inches(0.9), Inches(4.6), Inches(3), Inches(0.3),
                 "MILESTONES", font_size=11, color=PURPLE, bold=True)

    milestones = [
        ("Week 1-2", "DONE", "Core engine, CIS/NIST/CCM\nrule mapping, Terraform +\ncontainer scanning, PDF reports", GREEN),
        ("Week 3-4", "DONE", "Live AWS scanning, topology\nvisualization, dashboard\nwith trend analysis", GREEN),
        ("Month 2", "IN PROGRESS", "Azure + GCP live parity,\nGitHub CI/CD, scheduled\nscans", ACCENT),
        ("Month 3", "PLANNED", "Enterprise pilot, RBAC,\nAPI gateway, SaaS\ndeployment", BLUE),
    ]
    for i, (period, status, desc, color) in enumerate(milestones):
        x = Inches(0.9) + i * Inches(2.9)
        # Status badge
        badge_w = Inches(1.5) if status == "IN PROGRESS" else Inches(0.9)
        badge = add_rounded_rect(slide, x, Inches(5.0), badge_w, Inches(0.3),
                                 fill_color=RGBColor(0x0D, 0x1A, 0x0D),
                                 border_color=color)
        add_text_box(slide, x + Inches(0.05), Inches(5.02), badge_w, Inches(0.25),
                     status, font_size=9, color=color, bold=True, alignment=PP_ALIGN.CENTER)

        add_text_box(slide, x, Inches(5.4), Inches(2.5), Inches(0.3),
                     period, font_size=15, color=color, bold=True)
        add_text_box(slide, x, Inches(5.7), Inches(2.7), Inches(1.0),
                     desc, font_size=11, color=TEXT_SEC)

    add_slide_number(slide, 6)


def build_slide_7(prs):
    """PoC Screenshots"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_logo_badge(slide, Inches(0.6), Inches(0.4))

    add_text_box(slide, Inches(0.6), Inches(1.1), Inches(10), Inches(0.5),
                 "Working Proof of Concept — Live Application", font_size=28, color=WHITE, bold=True)

    # Badge
    badge = add_rounded_rect(slide, Inches(0.6), Inches(1.75), Inches(2.5), Inches(0.4),
                             fill_color=RGBColor(0x0D, 0x2A, 0x15), border_color=GREEN)
    add_text_box(slide, Inches(0.8), Inches(1.78), Inches(2.3), Inches(0.35),
                 "● LIVE DEMO SCREENSHOTS", font_size=10, color=GREEN, bold=True)

    # Row 1 — 3 screenshots
    screenshots_row1 = [
        ("login", "Login — Enterprise Auth"),
        ("dashboard", "Dashboard — Compliance Overview"),
        ("monitoring", "AWS Monitoring — 180 Resources"),
    ]

    for i, (key, caption) in enumerate(screenshots_row1):
        x = Inches(0.6) + i * Inches(4.1)
        y = Inches(2.4)
        w = Inches(3.8)
        h = Inches(2.2)

        # Card background
        card = add_rounded_rect(slide, x, y, w, h + Inches(0.4),
                                fill_color=BG_CARD, border_color=RGBColor(0x2A, 0x2D, 0x3E))

        # Screenshot image
        img_path = SCREENSHOTS.get(key)
        if img_path and os.path.exists(img_path):
            try:
                slide.shapes.add_picture(img_path, x + Inches(0.1), y + Inches(0.1),
                                        w - Inches(0.2), h - Inches(0.2))
            except Exception as e:
                add_text_box(slide, x + Inches(0.1), y + Inches(0.5), w - Inches(0.2), Inches(0.5),
                             f"[{key} screenshot]", font_size=14, color=TEXT_SEC, alignment=PP_ALIGN.CENTER)

        add_text_box(slide, x + Inches(0.1), y + h + Inches(0.05), w - Inches(0.2), Inches(0.3),
                     caption, font_size=10, color=ACCENT, bold=True, alignment=PP_ALIGN.CENTER)

    # Row 2 — 2 screenshots
    screenshots_row2 = [
        ("topology", "Infrastructure Topology — Multi-Region Heatmap"),
        ("compliance", "Compliance Rules — 108 Controls (CIS/NIST/CCM)"),
    ]

    for i, (key, caption) in enumerate(screenshots_row2):
        x = Inches(0.6) + i * Inches(6.2)
        y = Inches(5.15)
        w = Inches(5.8)
        h = Inches(1.8)

        card = add_rounded_rect(slide, x, y, w, h + Inches(0.4),
                                fill_color=BG_CARD, border_color=RGBColor(0x2A, 0x2D, 0x3E))

        img_path = SCREENSHOTS.get(key)
        if img_path and os.path.exists(img_path):
            try:
                slide.shapes.add_picture(img_path, x + Inches(0.1), y + Inches(0.1),
                                        w - Inches(0.2), h - Inches(0.15))
            except Exception as e:
                add_text_box(slide, x + Inches(0.1), y + Inches(0.5), w - Inches(0.2), Inches(0.5),
                             f"[{key} screenshot]", font_size=14, color=TEXT_SEC, alignment=PP_ALIGN.CENTER)

        add_text_box(slide, x + Inches(0.1), y + h + Inches(0.05), w - Inches(0.2), Inches(0.3),
                     caption, font_size=10, color=ACCENT, bold=True, alignment=PP_ALIGN.CENTER)

    add_slide_number(slide, 7)


def build_slide_8(prs):
    """The Ask & Conclusion"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_logo_badge(slide, Inches(0.6), Inches(0.4))

    add_text_box(slide, Inches(0.6), Inches(1.1), Inches(10), Inches(0.5),
                 "The Ask & Conclusion", font_size=28, color=WHITE, bold=True)

    # Resources Needed
    card1 = add_rounded_rect(slide, Inches(0.6), Inches(1.8), Inches(5.7), Inches(2.3),
                             fill_color=BG_CARD, border_color=RGBColor(0x2A, 0x2D, 0x3E))
    add_text_box(slide, Inches(0.9), Inches(1.9), Inches(3), Inches(0.3),
                 "RESOURCES NEEDED", font_size=11, color=ACCENT, bold=True)
    add_bullet_text(slide, Inches(0.9), Inches(2.25), Inches(5.2), Inches(1.7), [
        "GPU compute for Cisco Sec-8B model hosting & inference",
        "Azure & GCP sandbox accounts for live validation",
        "Security mentorship for adversarial AI testing",
        "Enterprise pilot partner for real-world validation",
    ], font_size=13)

    # TL;DR - Why we should win
    card2 = add_rounded_rect(slide, Inches(6.5), Inches(1.8), Inches(5.7), Inches(2.3),
                             fill_color=RGBColor(0x1A, 0x15, 0x05),
                             border_color=ACCENT)

    # Orange top bar on this card
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.5), Inches(1.8), Inches(5.7), Pt(3))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()

    add_text_box(slide, Inches(6.8), Inches(1.95), Inches(5), Inches(0.3),
                 "WHY THIS PROJECT SHOULD WIN", font_size=12, color=ACCENT, bold=True)

    win_reasons = [
        "First unified multi-framework compliance engine — CIS + NIST + CCM in a single AI-powered scan across AWS, Azure, GCP, Terraform, Docker, and Kubernetes",
        "Privacy-first AI — Cisco Sec-8B runs entirely on private infrastructure; zero data leaves the org",
        "Working PoC with real results — 180+ live AWS resources scanned, 63 critical findings identified, downloadable compliance reports generated",
    ]

    y = Inches(2.35)
    for i, reason in enumerate(win_reasons):
        # Number badge
        num_badge = add_rounded_rect(slide, Inches(6.8), y, Inches(0.3), Inches(0.3),
                                     fill_color=ACCENT)
        add_text_box(slide, Inches(6.8), y + Inches(0.02), Inches(0.3), Inches(0.25),
                     str(i + 1), font_size=12, color=BG_DARK, bold=True, alignment=PP_ALIGN.CENTER)
        
        add_text_box(slide, Inches(7.2), y, Inches(4.8), Inches(0.5),
                     reason, font_size=11, color=TEXT_SEC)
        y += Inches(0.6)

    # Bottom — Closing statement
    card3 = add_rounded_rect(slide, Inches(0.6), Inches(4.5), Inches(11.6), Inches(2.3),
                             fill_color=BG_CARD, border_color=RGBColor(0x2A, 0x2D, 0x3E))

    add_text_box(slide, Inches(2.5), Inches(4.9), Inches(8), Inches(0.5),
                 "Invecto Compliance Guard", font_size=30, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(2.5), Inches(5.5), Inches(8), Inches(0.4),
                 "AI-Powered Cloud Compliance — Simplified.", font_size=18, color=ACCENT,
                 bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(2.5), Inches(6.0), Inches(8), Inches(0.4),
                 "Powered by Cisco Sec-8B  ·  Invecto Technologies", font_size=12, color=TEXT_SEC,
                 alignment=PP_ALIGN.CENTER)

    add_slide_number(slide, 8)


# ═══════════════════════════════════════════════════════════════════════════════
# MAIN
# ═══════════════════════════════════════════════════════════════════════════════

def main():
    prs = Presentation()
    # Set to widescreen 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    print("Building Slide 1: Title & Team Overview...")
    build_slide_1(prs)

    print("Building Slide 2: Problem Statement...")
    build_slide_2(prs)

    print("Building Slide 3: Proposed AI Solution...")
    build_slide_3(prs)

    print("Building Slide 4: Technical Innovation...")
    build_slide_4(prs)

    print("Building Slide 5: Business Impact...")
    build_slide_5(prs)

    print("Building Slide 6: Roadmap...")
    build_slide_6(prs)

    print("Building Slide 7: PoC Screenshots...")
    build_slide_7(prs)

    print("Building Slide 8: The Ask & Conclusion...")
    build_slide_8(prs)

    prs.save(OUTPUT)
    print(f"\n✅ Presentation saved to: {OUTPUT}")
    print(f"   Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
